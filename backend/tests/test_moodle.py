"""Tests fuer Moodle-Endpunkte und Overview-Transformation."""

import io

import pytest

from app.services import moodle_service


def test_moodle_status_not_configured(client):
    response = client.get("/api/moodle/status")
    assert response.status_code == 200
    assert response.json() == {"configured": False}


def test_moodle_courses_without_token_returns_503(client):
    response = client.get("/api/moodle/courses")
    assert response.status_code == 503


def test_moodle_course_content_without_token_returns_503(client):
    response = client.get("/api/moodle/course/123/content")
    assert response.status_code == 503
    assert "MOODLE_TOKEN" in response.json()["detail"]


def test_moodle_course_overview_without_token_returns_503(client):
    response = client.get("/api/moodle/course/123/overview")
    assert response.status_code == 503
    assert "MOODLE_TOKEN" in response.json()["detail"]


def test_moodle_course_grades_without_token_returns_503(client):
    response = client.get("/api/moodle/course/123/grades")
    assert response.status_code == 503
    assert "MOODLE_TOKEN" in response.json()["detail"]


def test_moodle_course_deadlines_without_token_returns_503(client):
    response = client.get("/api/moodle/course/123/deadlines")
    assert response.status_code == 503
    assert "MOODLE_TOKEN" in response.json()["detail"]


@pytest.mark.asyncio
async def test_moodle_courses_transform(monkeypatch):
    async def fake_user_id():
        return 7

    async def fake_call(wsfunction, **params):
        assert wsfunction == "core_enrol_get_users_courses"
        assert params["userid"] == 7
        return [
            {
                "id": 123,
                "fullname": "Webtechnologien",
                "shortname": "B3.1 WebTech",
                "visible": True,
                "startdate": 0,
            }
        ]

    monkeypatch.setattr(moodle_service, "get_moodle_user_id", fake_user_id)
    monkeypatch.setattr(moodle_service, "_call", fake_call)

    courses = await moodle_service.get_moodle_courses()

    assert courses == [
        {
            "id": 123,
            "shortname": "B3.1 WebTech",
            "fullname": "Webtechnologien",
            "visible": True,
            "semester": "Unbekanntes Semester",
            "startdate": 0,
        }
    ]


@pytest.mark.asyncio
async def test_moodle_course_overview_transforms_sections(monkeypatch):
    monkeypatch.setattr(moodle_service.settings, "moodle_token", "test-token")

    async def fake_content(_course_id):
        return [
            {
                "name": "Session 1",
                "modules": [
                    {
                        "name": "Slides Session 1",
                        "modname": "resource",
                        "contents": [
                            {
                                "type": "file",
                                "filename": "slides.pdf",
                                "fileurl": "https://example.test/slides.pdf?forcedownload=1",
                                "mimetype": "application/pdf",
                                "filesize": 2048,
                            }
                        ],
                    },
                    {
                        "name": "Python Exercises",
                        "modname": "url",
                        "url": "https://example.test/exercises",
                    },
                    {
                        "name": "Submission Presentation 1",
                        "modname": "assign",
                        "dates": [{"label": "Due date", "timestamp": 1893456000}],
                    },
                ],
            }
        ]

    monkeypatch.setattr(moodle_service, "get_moodle_course_content", fake_content)

    overview = await moodle_service.get_moodle_course_overview("42")

    assert overview[0]["section_name"] == "Session 1"
    assert overview[0]["items"][0]["type"] == "file"
    assert overview[0]["items"][0]["filename"] == "slides.pdf"
    assert overview[0]["items"][0]["open_url"] == "https://example.test/slides.pdf?forcedownload=1&token=test-token"
    assert overview[0]["items"][1]["type"] == "url"
    assert overview[0]["items"][2]["type"] == "assignment"
    assert overview[0]["items"][2]["due_date"]


@pytest.mark.asyncio
async def test_moodle_course_grades_transform(monkeypatch):
    async def fake_user_id():
        return 7

    async def fake_call(wsfunction, **params):
        assert wsfunction == "gradereport_user_get_grade_items"
        assert params == {"courseid": "123", "userid": 7}
        return {
            "usergrades": [
                {
                    "gradeitems": [
                        {
                            "itemname": "Quiz 1",
                            "itemmodule": "quiz",
                            "gradeformatted": "1.7",
                            "grademaxformatted": "100",
                            "percentageformatted": "85%",
                            "feedback": "<p>Good work</p>",
                        }
                    ]
                }
            ]
        }

    monkeypatch.setattr(moodle_service, "get_moodle_user_id", fake_user_id)
    monkeypatch.setattr(moodle_service, "_call", fake_call)

    result = await moodle_service.get_moodle_course_grades("123")

    assert result == {
        "course_id": 123,
        "grades": [
            {
                "name": "Quiz 1",
                "type": "quiz",
                "grade": "1.7",
                "max_grade": "100",
                "percentage": "85%",
                "feedback": "Good work",
            }
        ],
    }


@pytest.mark.asyncio
async def test_moodle_course_deadlines_transform(monkeypatch):
    async def fake_content(_course_id):
        return [
            {
                "name": "Session 1",
                "modules": [
                    {
                        "name": "Submission Presentation 1",
                        "modname": "assign",
                        "url": "https://example.test/assign/view.php?id=99",
                        "dates": [{"label": "Due date", "timestamp": 1893455940}],
                    },
                    {
                        "name": "Regular Forum",
                        "modname": "forum",
                        "dates": [{"label": "Due date", "timestamp": 1893455940}],
                    },
                ],
            }
        ]

    monkeypatch.setattr(moodle_service, "get_moodle_course_content", fake_content)

    result = await moodle_service.get_moodle_course_deadlines("123")

    assert result["course_id"] == 123
    assert result["deadlines"] == [
        {
            "name": "Submission Presentation 1",
            "type": "assignment",
            "course_id": 123,
            "section_name": "Session 1",
            "due_date": "2029-12-31T23:59:00+00:00",
            "status": "open",
            "url": "https://example.test/assign/view.php?id=99",
            "open_url": "https://example.test/assign/view.php?id=99",
        }
    ]


@pytest.mark.asyncio
async def test_moodle_course_deadlines_empty(monkeypatch):
    async def fake_content(_course_id):
        return [{"name": "Session 1", "modules": [{"name": "Slides", "modname": "resource"}]}]

    monkeypatch.setattr(moodle_service, "get_moodle_course_content", fake_content)

    result = await moodle_service.get_moodle_course_deadlines("123")

    assert result == {"course_id": 123, "deadlines": []}


def test_add_token_to_url_handles_urls_without_query():
    result = moodle_service.add_token_to_url("https://example.test/file.pdf", "abc")
    assert result == "https://example.test/file.pdf?token=abc"


def test_add_token_to_url_handles_urls_with_query():
    result = moodle_service.add_token_to_url(
        "https://example.test/file.pdf?forcedownload=1", "abc"
    )
    assert result == "https://example.test/file.pdf?forcedownload=1&token=abc"


@pytest.mark.asyncio
async def test_moodle_course_files_include_pptx_and_docx(monkeypatch):
    async def fake_call(wsfunction, **params):
        assert wsfunction == "core_course_get_contents"
        assert params["courseid"] == 123
        return [
            {
                "name": "Session 1",
                "modules": [
                    {
                        "name": "Slides",
                        "contents": [
                            {
                                "type": "file",
                                "filename": "slides.pptx",
                                "fileurl": "https://example.test/slides.pptx",
                            },
                            {
                                "type": "file",
                                "filename": "notes.docx",
                                "fileurl": "https://example.test/notes.docx",
                            },
                            {
                                "type": "file",
                                "filename": "archive.zip",
                                "fileurl": "https://example.test/archive.zip",
                            },
                        ],
                    }
                ],
            }
        ]

    monkeypatch.setattr(moodle_service, "_call", fake_call)

    files = await moodle_service.get_course_files(123)

    assert [f["filename"] for f in files] == ["slides.pptx", "notes.docx"]


def test_pptx_extraction_preserves_every_slide_title_and_content():
    from pptx import Presentation

    presentation = Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "Promises Grundlagen"
    slide1.placeholders[1].text = "Asynchrone Werte\nthen und catch"

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Async Await"
    slide2.placeholders[1].text = "Lesbarer Kontrollfluss\nFehlerbehandlung"

    buffer = io.BytesIO()
    presentation.save(buffer)

    text = moodle_service._extract_pptx_text(buffer.getvalue())

    assert "Slide 1: Promises Grundlagen" in text
    assert "Asynchrone Werte" in text
    assert "then und catch" in text
    assert "Slide 2: Async Await" in text
    assert "Lesbarer Kontrollfluss" in text
    assert "Fehlerbehandlung" in text
