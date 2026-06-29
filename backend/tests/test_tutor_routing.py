"""Tests for TutorAgent fast path vs. material/RAG path."""

from unittest.mock import AsyncMock, MagicMock

import pytest
from langchain_core.messages import AIMessage

from app.agents import tutor_agent


class _FakeDirectLlm:
    def __init__(self, answer: str):
        self.answer = answer
        self.calls = []

    async def ainvoke(self, messages):
        self.calls.append(messages)
        return AIMessage(content=self.answer)


@pytest.mark.asyncio
async def test_general_programming_question_uses_direct_tutor_path(monkeypatch):
    fake_llm = _FakeDirectLlm("Promises sind Platzhalter für asynchrone Ergebnisse.")
    create_agent = MagicMock()
    retrieve_context = AsyncMock()

    monkeypatch.setattr(tutor_agent, "get_llm", lambda temperature=0.0: fake_llm)
    monkeypatch.setattr(tutor_agent, "create_tutor_agent", create_agent)
    monkeypatch.setattr(tutor_agent, "retrieve_context", retrieve_context)

    result = await tutor_agent.run_tutor_agent("Erklär mir Promises in JavaScript.", MagicMock())

    assert "Promises" in result
    assert fake_llm.calls
    create_agent.assert_not_called()
    retrieve_context.assert_not_called()


@pytest.mark.asyncio
async def test_document_question_keeps_material_rag_path(monkeypatch):
    fake_agent = MagicMock()
    fake_agent.ainvoke = AsyncMock(return_value={"messages": [AIMessage(content="Aus dem PDF...")]})
    create_agent = MagicMock(return_value=fake_agent)
    direct_llm = MagicMock()

    monkeypatch.setattr(tutor_agent, "get_llm", direct_llm)
    monkeypatch.setattr(tutor_agent, "create_tutor_agent", create_agent)

    result = await tutor_agent.run_tutor_agent("Erklär mir dieses PDF.", MagicMock())

    assert result == "Aus dem PDF..."
    create_agent.assert_called_once()
    direct_llm.assert_not_called()


@pytest.mark.asyncio
async def test_selected_moodle_session_uses_filtered_rag_path(monkeypatch):
    fake_llm = _FakeDirectLlm("Session 1 erklaert aus Moodle-Kontext.")
    retrieve_context = AsyncMock(return_value="[Quelle 1 - session1.pdf]\nKontext")
    create_agent = MagicMock()

    monkeypatch.setattr(tutor_agent, "get_llm", lambda temperature=0.0: fake_llm)
    monkeypatch.setattr(tutor_agent, "retrieve_context", retrieve_context)
    monkeypatch.setattr(tutor_agent, "create_tutor_agent", create_agent)

    moodle_context = {
        "course_id": 58776,
        "course_name": "B5.3 Unternehmenssoftware",
        "sections": [
            {
                "section_name": "Session 1",
                "items": [
                    {
                        "name": "Session 1 Slides",
                        "filename": "session1.pdf",
                        "type": "file",
                    }
                ],
            }
        ],
    }

    result = await tutor_agent.run_tutor_agent(
        "Explique-moi Session 1",
        MagicMock(),
        user_id="local",
        moodle_context=moodle_context,
    )

    assert "Session 1" in result
    create_agent.assert_not_called()
    retrieve_context.assert_awaited_once()
    kwargs = retrieve_context.call_args.kwargs
    assert kwargs["source_filter"] == ["session1.pdf"]
    assert kwargs["metadata_filter"] == {"moodle": "1", "course_id": 58776}
    assert fake_llm.calls


@pytest.mark.asyncio
async def test_selected_moodle_material_indexes_on_demand_when_missing(monkeypatch):
    fake_llm = _FakeDirectLlm("PPTX aus Moodle erklaert.")
    retrieve_context = AsyncMock(side_effect=["", "[Quelle 1 - slides.pptx]\nPromise slides"])
    download_file_text = AsyncMock(return_value="Promise slides")
    index_text = MagicMock(return_value=3)

    monkeypatch.setattr(tutor_agent, "retrieve_context", retrieve_context)
    monkeypatch.setattr(tutor_agent, "has_indexed_source", lambda *args, **kwargs: False)
    monkeypatch.setattr(tutor_agent, "get_llm", lambda temperature=0.0: fake_llm)
    monkeypatch.setattr("app.services.moodle_service.download_file_text", download_file_text)
    monkeypatch.setattr("app.rag.pipeline.index_text", index_text)

    result = await tutor_agent.run_tutor_agent(
        "Explique slides.pptx",
        MagicMock(),
        chat_id="chat-1",
        user_id="local",
        moodle_context={
            "course_id": "58776",
            "course_name": "B5.3 Unternehmenssoftware",
            "sections": [
                {
                    "section_name": "Session 1",
                    "items": [
                        {
                            "name": "USW Session 01",
                            "filename": "slides.pptx",
                            "fileurl": "https://moodle.test/slides.pptx",
                        }
                    ],
                }
            ],
        },
    )

    assert "PPTX aus Moodle" in result
    assert retrieve_context.await_count == 2
    download_file_text.assert_awaited_once_with("https://moodle.test/slides.pptx", "slides.pptx")
    index_text.assert_called_once()
    assert index_text.call_args.args[0] == "slides.pptx"
    assert index_text.call_args.args[2] == "chat-1"
    assert index_text.call_args.args[3] == "local"
    assert index_text.call_args.args[4]["moodle"] == "1"
    assert index_text.call_args.args[4]["course_id"] == 58776


@pytest.mark.asyncio
async def test_moodle_pptx_download_index_and_explain_integration(monkeypatch):
    from app.rag import pipeline, retriever
    from app.services import moodle_service
    from pptx import Presentation
    import io

    presentation = Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "Promise Grundlagen"
    slide1.placeholders[1].text = "Promises repraesentieren asynchrone Ergebnisse."
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Async Await"
    slide2.placeholders[1].text = "Await wartet auf die Erfuellung einer Promise."
    buffer = io.BytesIO()
    presentation.save(buffer)

    class _FakeResponse:
        content = buffer.getvalue()

        def raise_for_status(self):
            return None

    class _FakeClient:
        async def __aenter__(self):
            return self

        async def __aexit__(self, exc_type, exc, tb):
            return False

        async def get(self, url):
            assert "token=" in url
            return _FakeResponse()

    class _MemoryCollection:
        def __init__(self):
            self.documents = []
            self.metadatas = []

        def get(self, where=None, limit=None):
            matches = self._matches(where)
            return {"ids": [f"id-{i}" for i in matches[: limit or len(matches)]]}

        def upsert(self, ids, documents, embeddings, metadatas):
            self.documents.extend(documents)
            self.metadatas.extend(metadatas)

        def query(self, **kwargs):
            where = kwargs.get("where")
            matches = self._matches(where)
            docs = [self.documents[i] for i in matches]
            metas = [self.metadatas[i] for i in matches]
            return {"documents": [docs], "metadatas": [metas], "distances": [[0.1] * len(docs)]}

        def _matches(self, where):
            if not where:
                return list(range(len(self.documents)))

            def match_condition(meta, condition):
                if "$and" in condition:
                    return all(match_condition(meta, c) for c in condition["$and"])
                for key, value in condition.items():
                    if "$eq" in value and meta.get(key) != value["$eq"]:
                        return False
                    if "$in" in value and meta.get(key) not in value["$in"]:
                        return False
                return True

            return [i for i, meta in enumerate(self.metadatas) if match_condition(meta, where)]

    class _EchoLlm:
        async def ainvoke(self, messages):
            content = messages[-1].content
            assert "Promise Grundlagen" in content
            assert "Async Await" in content
            return AIMessage(content="Promises und Async/Await aus den Folien erklaert.")

    collection = _MemoryCollection()

    monkeypatch.setattr(moodle_service.settings, "moodle_token", "token")
    monkeypatch.setattr(moodle_service.httpx, "AsyncClient", lambda **_kwargs: _FakeClient())
    monkeypatch.setattr(pipeline, "get_collection", lambda: collection)
    monkeypatch.setattr(retriever, "get_collection", lambda: collection)
    monkeypatch.setattr(pipeline, "embed_texts", lambda texts: [[0.1] for _ in texts])
    monkeypatch.setattr(retriever, "embed_texts", lambda texts: [[0.1] for _ in texts])
    monkeypatch.setattr(tutor_agent, "get_llm", lambda temperature=0.0: _EchoLlm())

    result = await tutor_agent.run_tutor_agent(
        "Erkläre sample.pptx",
        MagicMock(),
        chat_id="chat-1",
        user_id="local",
        moodle_context={
            "course_id": 58776,
            "course_name": "B5.3 Unternehmenssoftware",
            "sections": [
                {
                    "section_name": "Session 1",
                    "items": [
                        {
                            "name": "Session 1 Upload",
                            "filename": "sample.pptx",
                            "fileurl": "https://moodle.test/sample.pptx",
                        }
                    ],
                }
            ],
        },
    )

    assert "Promises und Async/Await" in result
    assert any(meta.get("source") == "sample.pptx" for meta in collection.metadatas)
    assert any(meta.get("slide_number") == 1 for meta in collection.metadatas)
    assert any(meta.get("slide_number") == 2 for meta in collection.metadatas)
    assert all(meta.get("user_id") == "local" for meta in collection.metadatas)
    assert all(meta.get("chat_id") == "chat-1" for meta in collection.metadatas)
